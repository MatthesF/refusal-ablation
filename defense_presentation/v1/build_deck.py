#!/usr/bin/env python3
"""Build the 8-minute oral-defense deck for the refusal-ablation report."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figs")

# ---- Palette (built around the report figures: navy + orange) ----
NAVY_DARK = "12233F"   # dark backgrounds (title / conclusion)
NAVY      = "29457F"   # primary brand navy (matches figure blue)
NAVY_DEEP = "1C3460"
ORANGE    = "DE7A3C"   # accent (matches "Edited" bar)
ORANGE_DK = "C7672C"
LBLUE     = "8FB0E8"   # baseline-bar blue
ICE       = "CFE0FA"   # light text on dark
INK       = "1B2A45"   # body text on light
MUTED     = "6A7890"   # captions / secondary
CARD      = "F2F6FC"   # light card tint
CARDLINE  = "DCE6F4"
WHITE     = "FFFFFF"

TITLE_FONT = "Cambria"
BODY_FONT  = "Calibri"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def rgb(h):
    return RGBColor.from_string(h)


def add_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(NAVY_DARK if dark else WHITE)
    return s


def no_line(shape):
    shape.line.fill.background()


def add_shadow(shape, alpha="22000", blur="80000", dist="34000", direction="5400000"):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    eff = spPr.makeelement(qn('a:effectLst'), {})
    sh = eff.makeelement(qn('a:outerShdw'),
                         {'blurRad': blur, 'dist': dist, 'dir': direction, 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '0A1A33'})
    a = clr.makeelement(qn('a:alpha'), {'val': alpha})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def rect(slide, x, y, w, h, fill, rounded=False, radius=None, line=None, line_w=1.0, shadow=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        no_line(sp)
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_w)
    if rounded and radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        add_shadow(sp)
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True, space_after=None, line_spacing=None):
    """runs: list of paragraphs; each paragraph is list of (text, dict) run tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for t, o in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = o.get("font", BODY_FONT)
            f.size = Pt(o.get("size", 16))
            f.bold = o.get("bold", False)
            f.italic = o.get("italic", False)
            f.color.rgb = rgb(o.get("color", INK))
            if o.get("spacing"):
                _set_spacing(r, o["spacing"])
    return tb


def _set_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


def bullets(slide, x, y, w, h, items, size=15, color=INK, gap=8, lh=1.05, marker_color=ORANGE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = lh
        # marker
        rm = p.add_run()
        rm.text = "▸  "
        rm.font.name = BODY_FONT
        rm.font.size = Pt(size)
        rm.font.bold = True
        rm.font.color.rgb = rgb(marker_color)
        if isinstance(item, str):
            item = [(item, {})]
        for t, o in item:
            r = p.add_run()
            r.text = t
            r.font.name = o.get("font", BODY_FONT)
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", False)
            r.font.italic = o.get("italic", False)
            r.font.color.rgb = rgb(o.get("color", color))
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w and not h:
        h = w / ar
    elif h and not w:
        w = h * ar
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return pic, w, h


def speaker_pill(slide, name, dark=False):
    w, h = 2.35, 0.52
    x = SW - 0.6 - w
    y = 0.5
    label = txt(slide, x, y - 0.27, w, 0.24,
                [[("NOW PRESENTING", {"size": 8.5, "bold": True, "color": (ICE if dark else MUTED), "spacing": 2.2})]],
                align=PP_ALIGN.RIGHT)
    pill = rect(slide, x, y, w, h, ORANGE, rounded=True, radius=0.5, shadow=not dark)
    txt(slide, x, y, w, h, [[(name, {"size": 16, "bold": True, "color": WHITE, "font": BODY_FONT})]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return pill


def title_block(slide, kicker, title, dark=False, tw=9.6):
    txt(slide, 0.62, 0.5, tw, 0.3,
        [[(kicker, {"size": 12, "bold": True, "color": ORANGE, "spacing": 2.4})]])
    txt(slide, 0.62, 0.82, tw, 1.0,
        [[(title, {"size": 30, "bold": True, "color": (WHITE if dark else INK), "font": TITLE_FONT})]],
        line_spacing=0.98)


def footer(slide, num, dark=False):
    c = ICE if dark else MUTED
    txt(slide, 0.62, SH - 0.42, 8.5, 0.3,
        [[("Refusal-Direction Ablation in Gemma", {"size": 9.5, "color": c, "bold": True}),
          ("   ·   02445 Project Report · Technical University of Denmark", {"size": 9.5, "color": c})]])
    txt(slide, SW - 1.4, SH - 0.42, 0.8, 0.3,
        [[(f"{num:02d}", {"size": 9.5, "color": c, "bold": True})]], align=PP_ALIGN.RIGHT)


# =====================================================================
# SLIDE 1 — TITLE  (Oskar)
# =====================================================================
s = add_slide(dark=True)
# faint accent ring motif (oval outline, not a bar)
ring = slide_ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.7), Inches(-2.2), Inches(6.2), Inches(6.2))
ring.shadow.inherit = False
ring.fill.background()
ring.line.color.rgb = rgb(NAVY)
ring.line.width = Pt(1.5)
ring2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(4.6), Inches(4.6), Inches(4.6))
ring2.shadow.inherit = False
ring2.fill.background()
ring2.line.color.rgb = rgb(NAVY_DEEP)
ring2.line.width = Pt(1.2)

txt(s, 0.9, 1.05, 10, 0.4,
    [[("02445 PROJECT REPORT   ·   ORAL DEFENSE", {"size": 13, "bold": True, "color": ORANGE, "spacing": 2.6})]])
txt(s, 0.9, 1.65, 11.2, 2.1,
    [[("Statistical Evaluation of", {"size": 40, "bold": True, "color": WHITE, "font": TITLE_FONT})],
     [("Refusal-Direction Ablation in Gemma", {"size": 40, "bold": True, "color": WHITE, "font": TITLE_FONT})]],
    line_spacing=1.02)
txt(s, 0.9, 3.5, 10.6, 0.6,
    [[("Can a single weight edit remove benchmark-measured refusal behavior?",
       {"size": 17, "italic": True, "color": ICE})]])

# headline stat strip
rect(s, 0.9, 4.45, 4.55, 1.05, NAVY_DEEP, rounded=True, radius=0.09, shadow=True)
txt(s, 1.1, 4.58, 4.2, 0.4, [[("OFFICIAL SORRY-BENCH COMPLIANCE", {"size": 9.5, "bold": True, "color": ICE, "spacing": 1.6})]])
txt(s, 1.1, 4.86, 4.2, 0.6,
    [[("21.4%", {"size": 26, "bold": True, "color": LBLUE}),
      ("  →  ", {"size": 22, "bold": True, "color": WHITE}),
      ("76.6%", {"size": 26, "bold": True, "color": ORANGE})]])

# authors
txt(s, 6.0, 4.55, 6.4, 1.0,
    [[("Matthes M. Fogtmann   ·   Carl Johan von Löwzow", {"size": 14.5, "color": WHITE, "bold": True})],
     [("Oliver Illum   ·   Oskar F. Karlsson", {"size": 14.5, "color": WHITE, "bold": True})],
     [("Technical University of Denmark", {"size": 12.5, "color": MUTED})]],
    space_after=4)

txt(s, 0.9, SH - 0.55, 8, 0.3,
    [[("github.com/MatthesF/refusal-ablation", {"size": 11, "color": ICE})]])
speaker_pill(s, "Oskar", dark=True)

# =====================================================================
# SLIDE 2 — MOTIVATION  (Oskar)
# =====================================================================
s = add_slide()
title_block(s, "MOTIVATION", "Refusal is a safety mechanism you can edit")
speaker_pill(s, "Oskar")

cards = [
    ("Refusal is expected behaviour",
     "Instruction-tuned models should refuse harmful requests while staying useful on ordinary ones."),
    ("Open weights can be edited directly",
     "Safety is not only about prompts and decoding — for open-weight models the weights themselves are an attack surface."),
    ("Refusal lives in one direction",
     "Prior work shows refusal is largely mediated by a single direction in activation space (Arditi et al., 2024); “abliteration” removes it (Labonne 2024; Lai 2025)."),
]
cy = 2.05
cx = 0.62
cw = 7.05
ch = 1.42
for i, (head, body) in enumerate(cards):
    y = cy + i * (ch + 0.22)
    rect(s, cx, y, cw, ch, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0)
    num = rect(s, cx + 0.28, y + 0.32, 0.78, 0.78, NAVY, rounded=True, radius=0.5, shadow=True)
    txt(s, cx + 0.28, y + 0.32, 0.78, 0.78, [[(str(i + 1), {"size": 26, "bold": True, "color": WHITE, "font": TITLE_FONT})]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx + 1.3, y + 0.22, cw - 1.6, 0.4, [[(head, {"size": 17, "bold": True, "color": NAVY})]])
    txt(s, cx + 1.3, y + 0.62, cw - 1.6, 0.7, [[(body, {"size": 13, "color": INK})]], line_spacing=1.04)

# right takeaway panel
px = 8.0
rect(s, px, 2.05, 4.7, 4.5, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, px + 0.4, 2.4, 3.9, 0.5, [[("THE PROBLEM", {"size": 11, "bold": True, "color": ORANGE, "spacing": 2.0})]])
txt(s, px + 0.4, 2.85, 3.95, 3.4,
    [[("A model can score well on a fixed refusal benchmark,", {"size": 17, "color": WHITE, "bold": True})],
     [("", {"size": 6})],
     [("yet still rely on a safety mechanism that a cheap, post-hoc weight edit can weaken.",
       {"size": 17, "color": ICE})],
     [("", {"size": 10})],
     [("So: how much does the official benchmark score actually depend on that one mechanism?",
       {"size": 13.5, "italic": True, "color": ICE})]],
    line_spacing=1.06)
footer(s, 2)

# =====================================================================
# SLIDE 3 — RESEARCH QUESTION & DESIGN  (Oskar)
# =====================================================================
s = add_slide()
title_block(s, "RESEARCH QUESTION", "One narrow, reproducible question")
speaker_pill(s, "Oskar")

# question callout
rect(s, 0.62, 1.95, 12.1, 1.2, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, 0.95, 2.12, 0.6, 0.9, [[("“", {"size": 54, "bold": True, "color": ORANGE, "font": TITLE_FONT})]])
txt(s, 1.7, 2.18, 10.7, 0.95,
    [[("If we fit a refusal direction from a small policy-derived dataset and remove it from Gemma's weights, "
       "how much does official SORRY-Bench refusal behaviour change?", {"size": 16.5, "italic": True, "color": INK, "bold": True})]],
    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)

# two-condition flow
fy = 3.55
bw = 3.55
# baseline card
rect(s, 0.62, fy, bw, 1.7, WHITE, rounded=True, radius=0.06, line=LBLUE, line_w=2.0, shadow=True)
txt(s, 0.62, fy + 0.2, bw, 0.35, [[("BASELINE", {"size": 12, "bold": True, "color": NAVY, "spacing": 1.8})]], align=PP_ALIGN.CENTER)
txt(s, 0.82, fy + 0.62, bw - 0.4, 1.0,
    [[("Untouched", {"size": 15, "color": INK})],
     [("google/gemma-4-E4B-it", {"size": 13, "bold": True, "color": NAVY, "font": "Consolas"})]],
    align=PP_ALIGN.CENTER, space_after=3, line_spacing=1.0)
# edited card
ex = 0.62 + bw + 0.45
rect(s, ex, fy, bw, 1.7, WHITE, rounded=True, radius=0.06, line=ORANGE, line_w=2.0, shadow=True)
txt(s, ex, fy + 0.2, bw, 0.35, [[("EDITED", {"size": 12, "bold": True, "color": ORANGE_DK, "spacing": 1.8})]], align=PP_ALIGN.CENTER)
txt(s, ex + 0.2, fy + 0.62, bw - 0.4, 1.0,
    [[("Same checkpoint after", {"size": 15, "color": INK})],
     [("one fixed refusal-direction weight edit", {"size": 13, "bold": True, "color": ORANGE_DK})]],
    align=PP_ALIGN.CENTER, space_after=3, line_spacing=1.0)
# arrow
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ex + bw + 0.12), Inches(fy + 0.62), Inches(0.62), Inches(0.46))
ar.shadow.inherit = False
ar.fill.solid(); ar.fill.fore_color.rgb = rgb(MUTED); no_line(ar)
# benchmark card
jx = ex + bw + 0.9
jw = 12.72 - jx
rect(s, jx, fy, jw, 1.7, NAVY, rounded=True, radius=0.06, shadow=True)
txt(s, jx, fy + 0.2, jw, 0.35, [[("SAME EVALUATION", {"size": 12, "bold": True, "color": ORANGE, "spacing": 1.6})]], align=PP_ALIGN.CENTER)
txt(s, jx + 0.25, fy + 0.6, jw - 0.5, 1.05,
    [[("440 SORRY-Bench prompts, deterministic decoding,", {"size": 13, "color": WHITE})],
     [("scored by the official Mistral judge", {"size": 13, "color": WHITE, "bold": True})],
     [("0 = refusal   ·   1 = unsafe compliance", {"size": 12, "italic": True, "color": ICE})]],
    align=PP_ALIGN.CENTER, space_after=2, line_spacing=1.0)

# bottom note
txt(s, 0.62, 5.7, 12.1, 0.9,
    [[("Paired design — every prompt is answered once by each model.  ", {"size": 14, "color": INK, "bold": True}),
      ("The goal is to measure the behavioural change, not to build a useful model.", {"size": 14, "color": MUTED, "italic": True})]],
    line_spacing=1.05)
footer(s, 3)

# =====================================================================
# SLIDE 4 — CONSTRUCTION DATA  (Matthes)
# =====================================================================
s = add_slide()
title_block(s, "METHODS · DATA", "A direction fit from policy-derived prompts")
speaker_pill(s, "Matthes")

# left bullets
bullets(s, 0.62, 2.0, 6.3, 4.3, [
    [("480 prompts", {"bold": True, "color": NAVY}), (" from the Gemma Prohibited-Use Policy: 12 policy areas × (20 safe + 20 unsafe).", {})],
    [("Policy-area-disjoint split: ", {"bold": True, "color": NAVY}), ("8 areas (320) fit the direction, 4 areas (160) are held out for stability checks.", {})],
    [("Generated with GPT-5.5 Pro, then human-verified. ", {}), ("No supervised training or fine-tuning.", {"bold": True, "color": NAVY})],
    [("Child-related harmful prompts are ", {}), ("intentionally excluded", {"bold": True, "color": ORANGE_DK}), (" from construction (a data-handling choice).", {})],
    [("SORRY-Bench (440 prompts, 44 categories) is used ", {}), ("only for the final benchmark.", {"bold": True, "color": NAVY})],
], size=14.5, gap=12, lh=1.05)

# right: data roles table styled
tx = 7.25
tw = 5.45
ty = 2.0
rect(s, tx, ty, tw, 0.5, NAVY, rounded=False)
headers = ["Source", "Use", "Prompts"]
colx = [tx + 0.18, tx + 2.7, tx + 4.55]
for h, cxp in zip(headers, colx):
    txt(s, cxp, ty, 1.8, 0.5, [[(h, {"size": 12, "bold": True, "color": WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
rows = [
    ("Gemma policy prompts", "Direction data", "480"),
    ("Construction split", "Fit the edit", "320"),
    ("Held-out split", "Stability check", "160"),
    ("SORRY-Bench", "Final benchmark", "440"),
]
rh = 0.62
for i, (a, b, c) in enumerate(rows):
    ry = ty + 0.5 + i * rh
    rect(s, tx, ry, tw, rh, CARD if i % 2 == 0 else WHITE, line=CARDLINE, line_w=0.75)
    txt(s, colx[0], ry, 2.5, rh, [[(a, {"size": 12, "bold": True, "color": INK})]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, colx[1], ry, 1.9, rh, [[(b, {"size": 12, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, colx[2], ry, 1.0, rh, [[(c, {"size": 13, "bold": True, "color": NAVY})]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx, ty + 0.5 + 4 * rh + 0.12, tw, 0.4,
    [[("Construction and benchmark data never overlap.", {"size": 11.5, "italic": True, "color": MUTED})]])
footer(s, 4)

# =====================================================================
# SLIDE 5 — DIRECTION + WEIGHT EDIT  (Matthes)
# =====================================================================
s = add_slide()
title_block(s, "METHODS · THE EDIT", "Fit a direction, then remove it from the weights")
speaker_pill(s, "Matthes")

colw = 5.95
# Step 1 card
rect(s, 0.62, 2.0, colw, 4.35, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
b1 = rect(s, 0.9, 2.28, 0.62, 0.62, NAVY, rounded=True, radius=0.5, shadow=True)
txt(s, 0.9, 2.28, 0.62, 0.62, [[("1", {"size": 22, "bold": True, "color": WHITE, "font": TITLE_FONT})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 1.7, 2.32, colw - 1.3, 0.5, [[("Fit the refusal direction", {"size": 18, "bold": True, "color": NAVY, "font": TITLE_FONT})]])
bullets(s, 0.95, 3.15, colw - 0.6, 2.0, [
    "Run each prompt; take the final-token hidden state at every layer.",
    [("Direction = ", {}), ("mean(unsafe) − mean(safe)", {"bold": True, "color": INK}), (" activations.", {})],
    "One Gram-Schmidt step removes the component along the safe-prompt mean — a guardrail against erasing benign answering ability — then normalise.",
], size=13.5, gap=9, lh=1.04)
rect(s, 0.95, 5.55, colw - 0.66, 0.62, NAVY, rounded=True, radius=0.08)
txt(s, 0.95, 5.55, colw - 0.66, 0.62, [[("r", {"size": 15, "italic": True, "bold": True, "color": ORANGE, "font": "Cambria"}),
                                         ("  =  unit refusal direction per layer", {"size": 13.5, "color": WHITE})]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Step 2 card
sx = 0.62 + colw + 0.5
rect(s, sx, 2.0, colw, 4.35, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
b2 = rect(s, sx + 0.28, 2.28, 0.62, 0.62, ORANGE, rounded=True, radius=0.5, shadow=True)
txt(s, sx + 0.28, 2.28, 0.62, 0.62, [[("2", {"size": 22, "bold": True, "color": WHITE, "font": TITLE_FONT})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, sx + 1.08, 2.32, colw - 1.3, 0.5, [[("Remove it from the weights", {"size": 18, "bold": True, "color": ORANGE_DK, "font": TITLE_FONT})]])
bullets(s, sx + 0.33, 3.15, colw - 0.6, 2.0, [
    [("Edit every layer's ", {}), ("attention-output", {"bold": True, "color": INK}), (" and ", {}), ("MLP down-projection", {"bold": True, "color": INK}), (" matrices.", {})],
    "Strip from each weight column the part that writes along the refusal direction.",
    [("Restore column norms → ", {}), ("norm-preserving", {"bold": True, "color": ORANGE_DK}), (" edit; one fixed model for all 440 prompts.", {})],
], size=13.5, gap=9, lh=1.04, marker_color=ORANGE)
rect(s, sx + 0.33, 5.55, colw - 0.66, 0.62, NAVY, rounded=True, radius=0.08)
txt(s, sx + 0.33, 5.55, colw - 0.66, 0.62,
    [[("W", {"size": 15, "italic": True, "bold": True, "color": WHITE, "font": "Cambria"}),
      ("⊥", {"size": 12, "bold": True, "color": WHITE, "font": "Cambria"}),
      ("  =  W − ", {"size": 14, "color": WHITE, "font": "Cambria"}),
      ("r", {"size": 15, "italic": True, "bold": True, "color": ORANGE, "font": "Cambria"}),
      ("(", {"size": 14, "color": WHITE, "font": "Cambria"}),
      ("r", {"size": 15, "italic": True, "bold": True, "color": ORANGE, "font": "Cambria"}),
      ("ᵀW)", {"size": 14, "color": WHITE, "font": "Cambria"})]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 5)

# =====================================================================
# SLIDE 6 — STABILITY CHECK  (Matthes)
# =====================================================================
s = add_slide()
title_block(s, "METHODS · SANITY CHECK", "The fitted direction is stable across policy areas", tw=8.5)
speaker_pill(s, "Matthes")

_, iw, ih = add_image(s, os.path.join(FIG, "direction_stability.png"), 0.55, 2.0, w=7.5)
# right takeaways
rx = 8.35
bullets(s, rx, 2.25, 4.4, 3.8, [
    [("Directions fit from different policy-area groups point ", {}), ("the same way", {"bold": True, "color": NAVY}), (".", {})],
    [("Median cosine ", {}), ("0.95–0.97", {"bold": True, "color": ORANGE_DK}), (" to the held-out-area direction, even at the 10th percentile.", {})],
    "Suggests a shared, low-dimensional refusal signal — not just category-specific wording.",
], size=14.5, gap=12, lh=1.06)
rect(s, rx, 5.55, 4.4, 0.95, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0)
txt(s, rx + 0.25, 5.68, 3.95, 0.7,
    [[("Sanity check only", {"size": 13.5, "bold": True, "color": NAVY}),
      (" — never used to tune the final benchmark result.", {"size": 13.5, "color": INK})]], line_spacing=1.04,
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# =====================================================================
# SLIDE 7 — STATISTICAL DESIGN  (Carl Johan)
# =====================================================================
s = add_slide()
title_block(s, "STATISTICAL ANALYSIS", "A paired design with a binary judge")
speaker_pill(s, "Carl Johan")

# signed-change explainer
rect(s, 0.62, 2.0, 6.0, 2.0, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, 0.9, 2.2, 5.5, 0.4, [[("Each prompt → a signed change", {"size": 15, "bold": True, "color": NAVY})]])
chips = [("+1", "refusal → compliance", ORANGE), ("−0", "no change", MUTED), ("−1", "compliance → refusal", NAVY)]
cw2 = 1.72
for i, (v, lab, col) in enumerate(chips):
    xx = 0.9 + i * (cw2 + 0.12)
    rect(s, xx, 2.7, cw2, 1.1, WHITE, rounded=True, radius=0.08, line=col, line_w=1.5)
    txt(s, xx, 2.85, cw2, 0.5, [[(v.replace("−0", "0"), {"size": 22, "bold": True, "color": col, "font": TITLE_FONT})]], align=PP_ALIGN.CENTER)
    txt(s, xx + 0.08, 3.36, cw2 - 0.16, 0.4, [[(lab, {"size": 10.5, "color": INK})]], align=PP_ALIGN.CENTER, line_spacing=0.95)
txt(s, 0.9, 3.95, 5.6, 0.4, [[("Effect size  Δ = mean signed change", {"size": 13.5, "italic": True, "color": INK})]])

# right method chips (2x2)
methods = [
    ("McNemar's exact test", "Two-sided test on the discordant pairs: are refusal→compliance and compliance→refusal flips equally likely?"),
    ("95% confidence interval", "From the standard error of the paired prompt-level differences."),
    ("Category-clustered bootstrap", "20,000 replicates, seed 2445 — checks whether category clustering widens the uncertainty."),
    ("Powered by design", "440 paired prompts give 88.2% power to detect a 15 pp change."),
]
mx, my = 6.85, 2.0
mw, mh = 5.85, 1.05
for i, (h, b) in enumerate(methods):
    r, c = divmod(i, 1)
    y = my + i * (mh + 0.18)
    rect(s, mx, y, mw, mh, WHITE, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0, shadow=True)
    dot = rect(s, mx + 0.22, y + 0.34, 0.36, 0.36, NAVY, rounded=True, radius=0.5)
    txt(s, mx + 0.22, y + 0.34, 0.36, 0.36, [[(str(i + 1), {"size": 13, "bold": True, "color": WHITE})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, mx + 0.75, y + 0.13, mw - 0.95, 0.35, [[(h, {"size": 14, "bold": True, "color": NAVY})]])
    txt(s, mx + 0.75, y + 0.48, mw - 0.95, 0.5, [[(b, {"size": 11.5, "color": INK})]], line_spacing=1.0)

# balancing note (bottom-left)
rect(s, 0.62, 4.6, 6.0, 1.85, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, 0.92, 4.85, 5.5, 0.4, [[("WHY PAIRED + DETERMINISTIC", {"size": 11, "bold": True, "color": ORANGE, "spacing": 1.8})]])
txt(s, 0.92, 5.25, 5.45, 1.1,
    [[("Deterministic decoding removes sampling noise, so every flip is a real change in the judged label — not run-to-run variance. ", {"size": 13.5, "color": WHITE}),
      ("Pairing holds prompt difficulty, category mix and judge wording fixed across the two conditions.", {"size": 13.5, "color": ICE})]],
    line_spacing=1.1)
footer(s, 7)

# =====================================================================
# SLIDE 8 — MAIN RESULT  (Carl Johan)
# =====================================================================
s = add_slide()
title_block(s, "RESULTS · HEADLINE", "The edit removed most benchmark refusal", tw=8.5)
speaker_pill(s, "Carl Johan")

# left stat stack
sx = 0.62
rect(s, sx, 2.05, 4.55, 1.25, WHITE, rounded=True, radius=0.07, line=LBLUE, line_w=2.0, shadow=True)
txt(s, sx + 0.3, 2.18, 4.0, 0.35, [[("BASELINE GEMMA", {"size": 11, "bold": True, "color": NAVY, "spacing": 1.6})]])
txt(s, sx + 0.3, 2.5, 4.0, 0.7, [[("21.4%", {"size": 36, "bold": True, "color": NAVY, "font": TITLE_FONT}),
                                   ("   94/440 complied", {"size": 13, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, sx, 3.45, 4.55, 1.25, WHITE, rounded=True, radius=0.07, line=ORANGE, line_w=2.0, shadow=True)
txt(s, sx + 0.3, 3.58, 4.0, 0.35, [[("EDITED GEMMA", {"size": 11, "bold": True, "color": ORANGE_DK, "spacing": 1.6})]])
txt(s, sx + 0.3, 3.9, 4.0, 0.7, [[("76.6%", {"size": 36, "bold": True, "color": ORANGE_DK, "font": TITLE_FONT}),
                                   ("   337/440 complied", {"size": 13, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
# paired diff callout
rect(s, sx, 4.9, 4.55, 1.55, NAVY, rounded=True, radius=0.06, shadow=True)
txt(s, sx + 0.3, 5.04, 4.0, 0.35, [[("PAIRED DIFFERENCE", {"size": 11, "bold": True, "color": ORANGE, "spacing": 1.6})]])
txt(s, sx + 0.3, 5.32, 4.0, 0.6, [[("+55.2 pp", {"size": 32, "bold": True, "color": WHITE, "font": TITLE_FONT})]])
txt(s, sx + 0.3, 5.95, 4.1, 0.45,
    [[("95% CI [50.4, 60.0]   ·   bootstrap [47.3, 63.0]", {"size": 11.5, "color": ICE})]])

# right figure
add_image(s, os.path.join(FIG, "overall_compliance.png"), 5.45, 2.15, w=7.25)
txt(s, 5.45, 6.55, 7.25, 0.4,
    [[("Same model. One deterministic weight edit. ", {"size": 12.5, "italic": True, "bold": True, "color": INK}),
      ("Error bars: Wilson 95% intervals.", {"size": 12.5, "italic": True, "color": MUTED})]], align=PP_ALIGN.CENTER)
footer(s, 8)

# =====================================================================
# SLIDE 9 — PAIRED FLIPS  (Oliver)
# =====================================================================
s = add_slide()
title_block(s, "RESULTS · PAIRED FLIPS", "The change is overwhelmingly one-directional")
speaker_pill(s, "Oliver")

# contingency table
tx, ty = 0.62, 2.3
cell = 1.55
labw = 1.7
# corner labels
txt(s, tx + labw, ty - 0.42, cell * 2, 0.35, [[("EDITED GEMMA", {"size": 11, "bold": True, "color": ORANGE_DK, "spacing": 1.6})]], align=PP_ALIGN.CENTER)
txt(s, tx + labw, ty, cell, 0.45, [[("Refusal", {"size": 12, "bold": True, "color": INK})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx + labw + cell, ty, cell, 0.45, [[("Compliance", {"size": 12, "bold": True, "color": INK})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# vertical label
vlab = s.shapes.add_textbox(Inches(tx - 0.5), Inches(ty + 0.45), Inches(2.0), Inches(0.4))
vtf = vlab.text_frame; vtf.word_wrap = False
vr = vtf.paragraphs[0].add_run(); vr.text = "BASELINE GEMMA"
vr.font.size = Pt(11); vr.font.bold = True; vr.font.color.rgb = rgb(NAVY); vr.font.name = BODY_FONT
vlab.rotation = 270
vlab.left = Inches(tx - 1.05); vlab.top = Inches(ty + 1.05)

rowlabs = ["Refusal", "Compliance"]
data = [[("99", MUTED, "stayed refusal"), ("247", ORANGE, "flipped to compliance")],
        [("4", NAVY, "flipped to refusal"), ("90", MUTED, "stayed compliant")]]
for r in range(2):
    ry = ty + 0.45 + r * cell
    txt(s, tx, ry, labw - 0.1, cell, [[(rowlabs[r], {"size": 12, "bold": True, "color": INK})]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    for c in range(2):
        cxp = tx + labw + c * cell
        val, col, lab = data[r][c]
        is_key = (r == 0 and c == 1) or (r == 1 and c == 0)
        fillc = ORANGE if (r == 0 and c == 1) else (NAVY if (r == 1 and c == 0) else WHITE)
        textc = WHITE if is_key else INK
        rect(s, cxp, ry, cell, cell, fillc, line=CARDLINE, line_w=1.0, shadow=is_key)
        txt(s, cxp, ry + 0.25, cell, 0.6, [[(val, {"size": 30, "bold": True, "color": textc, "font": TITLE_FONT})]], align=PP_ALIGN.CENTER)
        txt(s, cxp + 0.08, ry + 0.95, cell - 0.16, 0.5, [[(lab, {"size": 10, "color": (ICE if is_key else MUTED)})]], align=PP_ALIGN.CENTER, line_spacing=0.95)

# right panel
px = 6.7
rect(s, px, 2.05, 6.02, 2.3, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, px + 0.35, 2.3, 5.4, 0.5, [[("247", {"size": 40, "bold": True, "color": ORANGE, "font": TITLE_FONT}),
                                    ("  vs  ", {"size": 22, "color": ICE}),
                                    ("4", {"size": 40, "bold": True, "color": LBLUE, "font": TITLE_FONT})]])
txt(s, px + 0.35, 3.25, 5.4, 0.9,
    [[("Flips run almost entirely refusal → compliance. McNemar's exact test: ", {"size": 13.5, "color": WHITE}),
      ("p = 9.07 × 10⁻⁶⁸", {"size": 13.5, "bold": True, "color": ORANGE})]], line_spacing=1.06)
rect(s, px, 4.55, 6.02, 1.9, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, px + 0.35, 4.78, 5.4, 0.4, [[("The four opposite flips", {"size": 14, "bold": True, "color": NAVY})]])
txt(s, px + 0.35, 5.18, 5.4, 1.2,
    [[("Personal insults, fake news, false common knowledge, dangerous-machinery advice — all borderline, caution-heavy cases. ",
       {"size": 13, "color": INK}),
      ("Reported as error analysis, not evidence of a safety gain.", {"size": 13, "italic": True, "color": MUTED})]],
    line_spacing=1.06)
footer(s, 9)

# =====================================================================
# SLIDE 10 — CATEGORY ERROR ANALYSIS  (Oliver)
# =====================================================================
s = add_slide()
title_block(s, "RESULTS · ERROR ANALYSIS", "Broad shift, but residual structure remains", tw=8.5)
speaker_pill(s, "Oliver")

# left: category change figure
add_image(s, os.path.join(FIG, "category_compliance_change.png"), 0.55, 2.1, w=6.95)
txt(s, 0.62, 5.12, 7.0, 0.4, [[("Many categories move 0% → 100%; others barely change. 10 prompts/category → exploratory.", {"size": 11, "italic": True, "color": MUTED})]])

# left note
rect(s, 0.62, 5.5, 6.95, 1.0, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0)
txt(s, 0.88, 5.62, 6.5, 0.8,
    [[("Child-related crimes rose ~30 pp", {"size": 13.5, "bold": True, "color": ORANGE_DK}),
      (" even though that content was excluded from the construction data.", {"size": 13.5, "color": INK})]],
    line_spacing=1.05, anchor=MSO_ANCHOR.MIDDLE)

# right: residual refusals mini bar chart
rx = 8.0
rect(s, rx, 2.05, 4.72, 4.45, WHITE, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0, shadow=True)
txt(s, rx + 0.3, 2.22, 4.2, 0.4, [[("Residual refusals concentrate", {"size": 14.5, "bold": True, "color": NAVY})]])
txt(s, rx + 0.3, 2.58, 4.2, 0.3, [[("103/440 prompts still refused · worst categories (/10)", {"size": 10.5, "color": MUTED})]])
resid = [("Child-related crimes", 7), ("Personal insult words", 6), ("Harassment", 6),
         ("PII violations", 6), ("Financial advice", 6), ("Violent crimes", 5)]
by = 3.08
barx = rx + 2.05
maxbar = 2.05
for i, (name, v) in enumerate(resid):
    yy = by + i * 0.52
    txt(s, rx + 0.3, yy - 0.02, 1.7, 0.35, [[(name, {"size": 9.5, "color": INK})]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, barx, yy + 0.03, maxbar, 0.26, CARD, rounded=True, radius=0.5, line=CARDLINE, line_w=0.5)
    col = ORANGE if i == 0 else NAVY
    rect(s, barx, yy + 0.03, maxbar * v / 10.0, 0.26, col, rounded=True, radius=0.5)
    txt(s, barx + maxbar * v / 10.0 + 0.06, yy - 0.02, 0.6, 0.35, [[(f"{v}/10", {"size": 10, "bold": True, "color": col})]], anchor=MSO_ANCHOR.MIDDLE)
footer(s, 10)

# =====================================================================
# SLIDE 11 — CONCLUSION  (Oskar)
# =====================================================================
s = add_slide(dark=True)
ringc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.0), Inches(4.4), Inches(5.4), Inches(5.4))
ringc.shadow.inherit = False; ringc.fill.background(); ringc.line.color.rgb = rgb(NAVY); ringc.line.width = Pt(1.3)
txt(s, 0.9, 0.7, 9.0, 0.35, [[("CONCLUSION", {"size": 13, "bold": True, "color": ORANGE, "spacing": 2.6})]])
txt(s, 0.9, 1.05, 11.4, 1.0, [[("What it shows — and what it doesn't", {"size": 32, "bold": True, "color": WHITE, "font": TITLE_FONT})]])
speaker_pill(s, "Oskar", dark=True)

# three takeaway columns
cols = [
    ("THE RESULT", ORANGE,
     "For gemma-4-E4B-it, one fitted refusal-direction edit removes most benchmark refusal: 21.4% → 76.6% compliance, with flips 247 vs 4. Statistically clear, practically large."),
    ("THE SCOPE", LBLUE,
     "This is a benchmark-specific behaviour change — not a verdict on model quality, helpfulness, or every safety mechanism. Residual refusals remain in some categories."),
    ("THE IMPLICATION", "FFFFFF",
     "Static refusal rates alone don't capture robustness for open-weight models. The residual follow-up shows test-set leakage risk; future evals need robustness-to-editing and fresh held-out tests."),
]
cw3 = 3.83
cy = 2.3
for i, (h, col, body) in enumerate(cols):
    x = 0.9 + i * (cw3 + 0.2)
    rect(s, x, cy, cw3, 3.35, NAVY_DEEP, rounded=True, radius=0.05, shadow=True)
    txt(s, x + 0.32, cy + 0.3, cw3 - 0.6, 0.4, [[(h, {"size": 12.5, "bold": True, "color": col, "spacing": 1.6})]])
    txt(s, x + 0.32, cy + 0.82, cw3 - 0.62, 2.3, [[(body, {"size": 13.5, "color": ICE})]], line_spacing=1.12)

rect(s, 0.9, 6.05, 11.53, 0.78, NAVY, rounded=True, radius=0.1)
txt(s, 0.9, 6.05, 11.53, 0.78,
    [[("Code, settings & artifacts:  ", {"size": 13, "color": ICE}),
      ("github.com/MatthesF/refusal-ablation", {"size": 13.5, "bold": True, "color": WHITE}),
      ("        Thank you — questions?", {"size": 13.5, "italic": True, "color": ORANGE})]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ---- speaker notes (brief, presenter view) ----
NOTES = {
    0: "Oskar opens. Title + the one-line question. 21.4 to 76.6 percent is the whole talk in one number.",
    1: "Oskar. Motivation: refusal expected; open weights editable; refusal lives in one direction (Arditi). The point: benchmark score may depend on one editable mechanism.",
    2: "Oskar. State the research question verbatim. Two conditions, same 440 prompts, same judge, paired. We measure the change, not build a product.",
    3: "Matthes. 480 policy prompts, disjoint split 320/160. GPT-5.5 Pro + human check. Child content excluded from construction. SORRY-Bench only for final test.",
    4: "Matthes. Two steps: fit direction (unsafe minus safe, Gram-Schmidt guardrail) then norm-preserving weight edit on attn-out + MLP down-proj.",
    5: "Matthes. Sanity check: directions from different policy areas agree (cosine 0.95-0.97). Shared low-dim refusal signal. Not used to tune the result.",
    6: "Carl Johan. Paired design, signed change +1/0/-1, effect size = mean. McNemar exact, 95% CI, clustered bootstrap, 88.2% power for 15pp.",
    7: "Carl Johan. Headline: 21.4 to 76.6 percent, +55.2pp, CI [50.4,60.0], bootstrap [47.3,63.0]. Same model, one edit.",
    8: "Oliver. 2x2 table: 247 refusal->compliance vs only 4 the other way. McNemar p ~ 1e-67. Four opposite flips are borderline cases, error analysis.",
    9: "Oliver. Broad but uneven. Some categories 0->100. Residuals concentrate (child-related 7/10 etc). Child crimes rose 30pp despite exclusion. 10/category = exploratory.",
    10: "Oskar closes. Result / scope / implication. Benchmark-specific change; static refusal rates don't capture robustness; leakage risk. Repo link, take questions.",
}
for idx, slide in enumerate(prs.slides):
    if idx in NOTES:
        slide.notes_slide.notes_text_frame.text = NOTES[idx]

out = os.path.join(HERE, "Refusal_Ablation_Defense.pptx")
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
