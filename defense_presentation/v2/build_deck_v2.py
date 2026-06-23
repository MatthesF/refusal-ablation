#!/usr/bin/env python3
"""Build the v2 oral-defense VIDEO PITCH deck (02445 closed-door seminar).

Targets the 5-7 minute video-pitch format, an equal 4-way speaker split, a
single core-conclusions closing slide, and explicit coverage of the course
learning objectives + rubric Steps 1-5. Appendix slides back the individual
questioning.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figs")

# ---- Palette (shared with v1) ----
NAVY_DARK = "12233F"
NAVY      = "29457F"
NAVY_DEEP = "1C3460"
ORANGE    = "DE7A3C"
ORANGE_DK = "C7672C"
LBLUE     = "8FB0E8"
ICE       = "CFE0FA"
INK       = "1B2A45"
MUTED     = "6A7890"
CARD      = "F2F6FC"
CARDLINE  = "DCE6F4"
GREEN     = "2E7D5B"
WHITE     = "FFFFFF"

TITLE_FONT = "Cambria"
BODY_FONT  = "Calibri"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def rgb(h):
    return RGBColor.from_string(h)


def add_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(NAVY_DARK if dark else WHITE)
    return s


def no_line(shape):
    shape.line.fill.background()


def add_shadow(shape, alpha="22000", blur="80000", dist="34000", direction="5400000"):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    eff = spPr.makeelement(qn('a:effectLst'), {})
    sh = eff.makeelement(qn('a:outerShdw'),
                         {'blurRad': blur, 'dist': dist, 'dir': direction, 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '0A1A33'})
    a = clr.makeelement(qn('a:alpha'), {'val': alpha})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def rect(slide, x, y, w, h, fill, rounded=False, radius=None, line=None, line_w=1.0, shadow=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        no_line(sp)
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_w)
    if rounded and radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        add_shadow(sp)
    return sp


def _set_spacing(run, pts):
    rPr = run._r.get_or_add_rPr()
    rPr.set('spc', str(int(pts * 100)))


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        wrap=True, space_after=None, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None:
            p.space_after = Pt(space_after)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        for t, o in para:
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = o.get("font", BODY_FONT)
            f.size = Pt(o.get("size", 16))
            f.bold = o.get("bold", False)
            f.italic = o.get("italic", False)
            f.color.rgb = rgb(o.get("color", INK))
            if o.get("spacing"):
                _set_spacing(r, o["spacing"])
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=INK, gap=8, lh=1.05, marker_color=ORANGE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = lh
        rm = p.add_run()
        rm.text = "▸  "
        rm.font.name = BODY_FONT
        rm.font.size = Pt(size)
        rm.font.bold = True
        rm.font.color.rgb = rgb(marker_color)
        if isinstance(item, str):
            item = [(item, {})]
        for t, o in item:
            r = p.add_run()
            r.text = t
            r.font.name = o.get("font", BODY_FONT)
            r.font.size = Pt(o.get("size", size))
            r.font.bold = o.get("bold", False)
            r.font.italic = o.get("italic", False)
            r.font.color.rgb = rgb(o.get("color", color))
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w and not h:
        h = w / ar
    elif h and not w:
        w = h * ar
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return w, h


def speaker_pill(slide, name, dark=False, sublabel="NOW PRESENTING"):
    w, h = 2.45, 0.52
    x = SW - 0.6 - w
    y = 0.5
    txt(slide, x, y - 0.27, w, 0.24,
        [[(sublabel, {"size": 8.5, "bold": True, "color": (ICE if dark else MUTED), "spacing": 2.2})]],
        align=PP_ALIGN.RIGHT)
    rect(slide, x, y, w, h, ORANGE, rounded=True, radius=0.5, shadow=not dark)
    txt(slide, x, y, w, h, [[(name, {"size": 15.5, "bold": True, "color": WHITE})]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def title_block(slide, kicker, title, dark=False, tw=9.4):
    txt(slide, 0.62, 0.5, tw, 0.3,
        [[(kicker, {"size": 12, "bold": True, "color": ORANGE, "spacing": 2.4})]])
    txt(slide, 0.62, 0.82, tw, 1.0,
        [[(title, {"size": 29, "bold": True, "color": (WHITE if dark else INK), "font": TITLE_FONT})]],
        line_spacing=0.98)


def chips(slide, x, y, labels, color=NAVY):
    """small competency tag chips in a row"""
    cx = x
    for lab in labels:
        w = 0.16 + 0.092 * len(lab)
        rect(slide, cx, y, w, 0.34, CARD, rounded=True, radius=0.5, line=CARDLINE, line_w=1.0)
        txt(slide, cx, y, w, 0.34, [[(lab, {"size": 9.5, "bold": True, "color": color})]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w + 0.12


def footer(slide, num, dark=False, tag=None):
    c = ICE if dark else MUTED
    left = "Refusal-Direction Ablation in Gemma"
    if tag:
        left = tag
    txt(slide, 0.62, SH - 0.42, 9.5, 0.3,
        [[(left, {"size": 9.5, "color": c, "bold": True}),
          ("   ·   02445 · Video pitch (5–7 min) · DTU", {"size": 9.5, "color": c})]])
    txt(slide, SW - 1.4, SH - 0.42, 0.8, 0.3,
        [[(num, {"size": 9.5, "color": c, "bold": True})]], align=PP_ALIGN.RIGHT)


def two_condition_flow(s, fy):
    bw = 3.55
    rect(s, 0.62, fy, bw, 1.5, WHITE, rounded=True, radius=0.06, line=LBLUE, line_w=2.0, shadow=True)
    txt(s, 0.62, fy + 0.16, bw, 0.32, [[("A · BASELINE", {"size": 11.5, "bold": True, "color": NAVY, "spacing": 1.4})]], align=PP_ALIGN.CENTER)
    txt(s, 0.82, fy + 0.55, bw - 0.4, 0.9,
        [[("Untouched", {"size": 14, "color": INK})],
         [("gemma-4-E4B-it", {"size": 12.5, "bold": True, "color": NAVY, "font": "Consolas"})]],
        align=PP_ALIGN.CENTER, space_after=2, line_spacing=1.0)
    ex = 0.62 + bw + 0.45
    rect(s, ex, fy, bw, 1.5, WHITE, rounded=True, radius=0.06, line=ORANGE, line_w=2.0, shadow=True)
    txt(s, ex, fy + 0.16, bw, 0.32, [[("B · EDITED", {"size": 11.5, "bold": True, "color": ORANGE_DK, "spacing": 1.4})]], align=PP_ALIGN.CENTER)
    txt(s, ex + 0.2, fy + 0.55, bw - 0.4, 0.9,
        [[("Same checkpoint after", {"size": 14, "color": INK})],
         [("one refusal-direction edit", {"size": 12.5, "bold": True, "color": ORANGE_DK})]],
        align=PP_ALIGN.CENTER, space_after=2, line_spacing=1.0)
    ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ex + bw + 0.12), Inches(fy + 0.52), Inches(0.62), Inches(0.46))
    ar.shadow.inherit = False
    ar.fill.solid(); ar.fill.fore_color.rgb = rgb(MUTED); no_line(ar)
    jx = ex + bw + 0.9
    jw = 12.72 - jx
    rect(s, jx, fy, jw, 1.5, NAVY, rounded=True, radius=0.06, shadow=True)
    txt(s, jx, fy + 0.16, jw, 0.32, [[("SAME PAIRED TEST", {"size": 11.5, "bold": True, "color": ORANGE, "spacing": 1.4})]], align=PP_ALIGN.CENTER)
    txt(s, jx + 0.25, fy + 0.52, jw - 0.5, 0.95,
        [[("440 SORRY-Bench prompts,", {"size": 12.5, "color": WHITE})],
         [("deterministic decoding", {"size": 12.5, "bold": True, "color": WHITE})]],
        align=PP_ALIGN.CENTER, space_after=1, line_spacing=1.0)


# =====================================================================
# 1 — TITLE  (Oskar)
# =====================================================================
s = add_slide(dark=True)
ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.7), Inches(-2.2), Inches(6.2), Inches(6.2))
ring.shadow.inherit = False; ring.fill.background(); ring.line.color.rgb = rgb(NAVY); ring.line.width = Pt(1.5)
ring2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(4.6), Inches(4.6), Inches(4.6))
ring2.shadow.inherit = False; ring2.fill.background(); ring2.line.color.rgb = rgb(NAVY_DEEP); ring2.line.width = Pt(1.2)

txt(s, 0.9, 1.0, 10.5, 0.4,
    [[("02445 · DESIGN & EVALUATE A GENAI SYSTEM · VIDEO PITCH", {"size": 12.5, "bold": True, "color": ORANGE, "spacing": 2.2})]])
txt(s, 0.9, 1.6, 11.4, 2.1,
    [[("Statistical Evaluation of", {"size": 39, "bold": True, "color": WHITE, "font": TITLE_FONT})],
     [("Refusal-Direction Ablation in Gemma", {"size": 39, "bold": True, "color": WHITE, "font": TITLE_FONT})]],
    line_spacing=1.02)
txt(s, 0.9, 3.42, 10.8, 0.6,
    [[("A paired A/B test of whether one weight edit removes benchmark-measured refusal", {"size": 16.5, "italic": True, "color": ICE})]])

rect(s, 0.9, 4.4, 4.55, 1.05, NAVY_DEEP, rounded=True, radius=0.09, shadow=True)
txt(s, 1.1, 4.53, 4.2, 0.4, [[("OFFICIAL SORRY-BENCH COMPLIANCE", {"size": 9.5, "bold": True, "color": ICE, "spacing": 1.6})]])
txt(s, 1.1, 4.81, 4.2, 0.6,
    [[("21.4%", {"size": 26, "bold": True, "color": LBLUE}),
      ("  →  ", {"size": 22, "bold": True, "color": WHITE}),
      ("76.6%", {"size": 26, "bold": True, "color": ORANGE})]])
txt(s, 6.0, 4.5, 6.4, 1.0,
    [[("Matthes M. Fogtmann   ·   Carl Johan von Löwzow", {"size": 14, "color": WHITE, "bold": True})],
     [("Oliver Illum   ·   Oskar F. Karlsson", {"size": 14, "color": WHITE, "bold": True})],
     [("Technical University of Denmark", {"size": 12, "color": MUTED})]],
    space_after=4)
txt(s, 0.9, SH - 0.55, 8, 0.3, [[("github.com/MatthesF/refusal-ablation", {"size": 11, "color": ICE})]])
speaker_pill(s, "Oskar", dark=True)

# =====================================================================
# 2 — STEP 1: ASPECT & WHY  (Oskar)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEP 1 · WHAT WE EVALUATE", "Robustness of an open-weight model's refusal")
speaker_pill(s, "Oskar")

cards = [
    ("The aspect", "Robustness / red-teaming: does a single, cheap weight edit remove the refusal behaviour a static safety benchmark measures?"),
    ("Why it matters", "Open-weight models can be edited after release. A high benchmark refusal rate is only meaningful if it survives such edits."),
    ("What's at stake", "If a post-hoc edit flips the score, the benchmark measures the prompt, not a robust safety mechanism — a validity problem for static evals."),
]
cy, cx, cw, ch = 2.0, 0.62, 7.05, 1.42
for i, (head, body) in enumerate(cards):
    y = cy + i * (ch + 0.22)
    rect(s, cx, y, cw, ch, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0)
    rect(s, cx + 0.28, y + 0.32, 0.78, 0.78, NAVY, rounded=True, radius=0.5, shadow=True)
    txt(s, cx + 0.28, y + 0.32, 0.78, 0.78, [[(str(i + 1), {"size": 26, "bold": True, "color": WHITE, "font": TITLE_FONT})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, cx + 1.3, y + 0.24, cw - 1.6, 0.4, [[(head, {"size": 17, "bold": True, "color": NAVY})]])
    txt(s, cx + 1.3, y + 0.64, cw - 1.6, 0.7, [[(body, {"size": 12.5, "color": INK})]], line_spacing=1.04)

px = 8.0
rect(s, px, 2.0, 4.7, 4.5, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, px + 0.4, 2.32, 3.9, 0.5, [[("OUR CHOICE", {"size": 11, "bold": True, "color": ORANGE, "spacing": 2.0})]])
txt(s, px + 0.4, 2.78, 3.95, 3.5,
    [[("We red-team Gemma's safety in weight space", {"size": 18, "color": WHITE, "bold": True})],
     [("", {"size": 6})],
     [("Not adversarial prompts — a fixed, norm-preserving edit fit from a policy-derived dataset.", {"size": 14, "color": ICE})],
     [("", {"size": 8})],
     [("Question: how much of SORRY-Bench refusal depends on the single direction we remove?", {"size": 13.5, "italic": True, "color": ICE})]],
    line_spacing=1.08)
footer(s, "02")

# =====================================================================
# 3 — DESIGN: A/B TEST + METRIC  (Oskar)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEPS 2–3 · DESIGN & METRIC", "A paired A/B test with an LLM-as-a-Judge metric", tw=10.5)
speaker_pill(s, "Oskar")

two_condition_flow(s, 1.95)

# metric card
rect(s, 0.62, 3.75, 7.0, 2.05, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, 0.92, 3.95, 6.4, 0.4, [[("Metric (Step 2): LLM-as-a-Judge", {"size": 15.5, "bold": True, "color": NAVY})]])
bullets(s, 0.92, 4.4, 6.45, 1.4, [
    [("Official SORRY-Bench fine-tuned ", {}), ("Mistral-7B judge", {"bold": True, "color": INK}), (" — the exact LLM-as-a-Judge path the brief suggests.", {})],
    [("Binary score per answer: ", {}), ("0 = refusal", {"bold": True, "color": NAVY}), (",  ", {}), ("1 = unsafe compliance", {"bold": True, "color": ORANGE_DK}), (".", {})],
    "Manual audit of 40 outputs: official labels reasonable for 36/40 (judge is useful but coarse).",
], size=12.5, gap=7, lh=1.04)

rect(s, 7.85, 3.75, 4.87, 2.05, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, 8.15, 3.95, 4.3, 0.4, [[("WHY “PAIRED A/B”", {"size": 11, "bold": True, "color": ORANGE, "spacing": 1.8})]])
txt(s, 8.15, 4.38, 4.32, 1.4,
    [[("Both arms answer the ", {"size": 13.5, "color": WHITE}),
      ("same 440 prompts", {"size": 13.5, "bold": True, "color": WHITE}),
      (". Pairing holds prompt difficulty, category mix and judge wording fixed, so each prompt is its own control.", {"size": 13.5, "color": ICE})]],
    line_spacing=1.1)
chips(s, 0.62, 6.05, ["A/B test", "within-subject design", "LLM-as-a-Judge"])
footer(s, "03")

# =====================================================================
# 4 — MODEL & RUNTIME  (Matthes)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEP 3 · MODEL & RUNTIME", "An open-weight model we can edit and run locally")
speaker_pill(s, "Matthes")

specs = [
    ("Model", "google/gemma-4-E4B-it", "Instruction-tuned, open-weight, pinned to one revision for reproducibility."),
    ("Why this one", "Weight-space access", "Direct edits are impossible through a closed, hosted API — open weights are required."),
    ("Runtime", "RunPod RTX 5090 GPU", "Run locally/offline; deterministic decoding, 4096-token cap. (the GPU-exploration angle)"),
    ("Independence", "Offline & deterministic", "No train-on-prompts feedback loop, so the 440 paired observations stay independent."),
]
gx, gy, gw, gh = 0.62, 2.05, 5.9, 1.95
for i, (tag, big, body) in enumerate(specs):
    r, c = divmod(i, 2)
    x = gx + c * (gw + 0.3)
    y = gy + r * (gh + 0.25)
    rect(s, x, y, gw, gh, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0, shadow=True)
    txt(s, x + 0.3, y + 0.22, gw - 0.6, 0.3, [[(tag.upper(), {"size": 10.5, "bold": True, "color": ORANGE_DK, "spacing": 1.6})]])
    txt(s, x + 0.3, y + 0.55, gw - 0.6, 0.5, [[(big, {"size": 18, "bold": True, "color": NAVY, "font": TITLE_FONT})]])
    txt(s, x + 0.3, y + 1.12, gw - 0.6, 0.7, [[(body, {"size": 12.5, "color": INK})]], line_spacing=1.05)
footer(s, "04")

# =====================================================================
# 5 — DATA  (Matthes)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEP 4 · DATA", "A benchmark plus a self-curated construction set", tw=10.5)
speaker_pill(s, "Matthes")

bullets(s, 0.62, 2.0, 6.3, 3.0, [
    [("Benchmark (external hold-out): ", {"bold": True, "color": NAVY}), ("SORRY-Bench main split — 440 prompts, 44 harm categories, 10 each.", {})],
    [("Self-curated construction set: ", {"bold": True, "color": NAVY}), ("480 prompts from the Gemma Prohibited-Use Policy.", {})],
    [("Factors of variation: ", {"bold": True, "color": ORANGE_DK}), ("12 policy areas × 2 labels (safe / unsafe), 20 each — a crossed design.", {})],
    "Generated with GPT-5.5 Pro, then human-verified; policy-area-disjoint 8 / 4 split. No fine-tuning.",
], size=13.5, gap=11, lh=1.05)

# power callout
rect(s, 0.62, 5.45, 6.3, 1.05, NAVY, rounded=True, radius=0.06, shadow=True)
txt(s, 0.92, 5.6, 5.8, 0.35, [[("SAMPLE SIZE (Step 4d)", {"size": 10.5, "bold": True, "color": ORANGE, "spacing": 1.6})]])
txt(s, 0.92, 5.92, 5.8, 0.5,
    [[("n ≈ 349", {"size": 16, "bold": True, "color": WHITE}),
      (" needed for a 15 pp change at 80% power → 440 prompts give ", {"size": 12.5, "color": ICE}),
      ("88.2% power", {"size": 13, "bold": True, "color": WHITE}), (".", {"size": 12.5, "color": ICE})]],
    line_spacing=1.05)

# data-roles table
tx, tw2, ty = 7.25, 5.45, 2.0
rect(s, tx, ty, tw2, 0.5, NAVY)
headers = ["Source", "Role", "Prompts"]
colx = [tx + 0.18, tx + 2.7, tx + 4.55]
for h, cxp in zip(headers, colx):
    txt(s, cxp, ty, 1.8, 0.5, [[(h, {"size": 12, "bold": True, "color": WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
rows = [("Construction split", "Fit the edit", "320"),
        ("Held-out split", "Generalization", "160"),
        ("SORRY-Bench", "Final A/B test", "440")]
rh = 0.66
for i, (a, b, c) in enumerate(rows):
    ry = ty + 0.5 + i * rh
    rect(s, tx, ry, tw2, rh, CARD if i % 2 == 0 else WHITE, line=CARDLINE, line_w=0.75)
    txt(s, colx[0], ry, 2.5, rh, [[(a, {"size": 12, "bold": True, "color": INK})]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, colx[1], ry, 1.9, rh, [[(b, {"size": 12, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, colx[2], ry, 1.0, rh, [[(c, {"size": 13, "bold": True, "color": NAVY})]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx, ty + 0.5 + 3 * rh + 0.14, tw2, 0.5,
    [[("Construction data and the benchmark never overlap; child-related content excluded from construction.", {"size": 11, "italic": True, "color": MUTED})]], line_spacing=1.04)
footer(s, "05")

# =====================================================================
# 6 — INTERVENTION + RESPONSES → DATA  (Matthes)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEP 4c · OUTPUTS → A PROBABILITY MODEL", "Turning generated answers into modellable data", tw=11)
speaker_pill(s, "Matthes")

# intervention strip (brief)
rect(s, 0.62, 2.0, 12.1, 1.35, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, 0.92, 2.16, 11.5, 0.35, [[("The intervention (one fixed edit — detail in appendix)", {"size": 14, "bold": True, "color": NAVY})]])
txt(s, 0.92, 2.55, 11.6, 0.7,
    [[("Fit a refusal direction  ", {"size": 13, "color": INK}),
      ("r = mean(unsafe) − mean(safe)", {"size": 13, "bold": True, "color": INK, "font": "Cambria"}),
      ("  (orthogonalised, unit norm), then edit every layer's weights  ", {"size": 13, "color": INK}),
      ("W⊥ = W − r(rᵀW)", {"size": 13, "bold": True, "color": ORANGE_DK, "font": "Cambria"}),
      ("  with column norms restored (norm-preserving).", {"size": 13, "color": INK})]],
    line_spacing=1.06)

# pipeline → Bernoulli
fy = 3.7
steps = [("Prompt i", LBLUE), ("Model answer", NAVY), ("Mistral judge", ORANGE), ("Yᵢ ∈ {0, 1}", NAVY_DEEP)]
bw = 2.55
for i, (lab, col) in enumerate(steps):
    x = 0.62 + i * (bw + 0.42)
    rect(s, x, fy, bw, 0.95, col, rounded=True, radius=0.1, shadow=True)
    lab_font = BODY_FONT if "∈" in lab else TITLE_FONT
    txt(s, x, fy, bw, 0.95, [[(lab, {"size": 15, "bold": True, "color": WHITE, "font": lab_font})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.06), Inches(fy + 0.32), Inches(0.32), Inches(0.3))
        ar.shadow.inherit = False; ar.fill.solid(); ar.fill.fore_color.rgb = rgb(MUTED); no_line(ar)

rect(s, 0.62, 5.05, 12.1, 1.45, NAVY, rounded=True, radius=0.05, shadow=True)
bullets(s, 0.95, 5.22, 11.5, 1.2, [
    [("Each judged answer is a ", {"color": WHITE}), ("Bernoulli outcome", {"bold": True, "color": WHITE}), (" Y ∈ {0,1}; deterministic decoding fixes each label (no sampling noise).", {"color": ICE})],
    [("Per prompt we keep the ", {"color": WHITE}), ("paired outcome (Y_baseline, Y_edited)", {"bold": True, "color": WHITE}), (" — the unit of analysis for the A/B test.", {"color": ICE})],
], size=13, gap=6, lh=1.05, marker_color=ORANGE)
footer(s, "06")

# =====================================================================
# 7 — STATISTICAL MODEL  (Carl Johan)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEP 5b · THE MODEL", "The evaluation written as a statistical model", tw=10.5)
speaker_pill(s, "Carl Johan")

# left: paired effect
rect(s, 0.62, 2.0, 5.8, 2.0, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, 0.92, 2.18, 5.3, 0.35, [[("Paired effect size", {"size": 15, "bold": True, "color": NAVY})]])
txt(s, 0.92, 2.62, 5.4, 0.5,
    [[("Dᵢ = Y", {"size": 16, "color": INK, "font": "Cambria"}),
      ("edited", {"size": 11, "color": INK, "font": "Cambria"}),
      (" − Y", {"size": 16, "color": INK, "font": "Cambria"}),
      ("baseline", {"size": 11, "color": INK, "font": "Cambria"}),
      ("  ∈ {−1, 0, +1}", {"size": 16, "bold": True, "color": ORANGE_DK, "font": "Cambria"})]])
txt(s, 0.92, 3.22, 5.4, 0.6,
    [[("Δ = mean(Dᵢ)", {"size": 16, "bold": True, "color": NAVY, "font": "Cambria"}),
      ("   — the compliance-rate change", {"size": 12.5, "color": MUTED})]])

# right: GLM / matrix lens
rect(s, 6.6, 2.0, 6.12, 2.0, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, 6.9, 2.18, 5.5, 0.35, [[("Generalized-linear / matrix lens", {"size": 15, "bold": True, "color": ORANGE})]])
txt(s, 6.9, 2.62, 5.6, 0.9,
    [[("logit P(Yᵢ=1) = β₀ + β₁·edited", {"size": 14.5, "bold": True, "color": WHITE, "font": "Cambria"})],
     [("                       + Σ γ_c·catᵢ + Σ δ_c·(edited×catᵢ)", {"size": 13, "color": ICE, "font": "Cambria"})]],
    line_spacing=1.15)
txt(s, 6.9, 3.5, 5.6, 0.45,
    [[("η = Xβ", {"size": 13.5, "italic": True, "bold": True, "color": WHITE, "font": "Cambria"}),
      ("  — design matrix X holds condition, category and their interaction.", {"size": 11.5, "color": ICE})]], line_spacing=1.0)

# bottom: what each term demonstrates
items = [("β₁", "Main edit effect", "the A/B treatment effect"),
         ("δ_c", "Interaction", "does the edit's effect differ by category?"),
         ("X", "Matrix formulation", "simple model in design-matrix form")]
iy = 4.35
iw = 3.93
for i, (sym, head, body) in enumerate(items):
    x = 0.62 + i * (iw + 0.18)
    rect(s, x, iy, iw, 1.75, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0)
    rect(s, x + 0.28, iy + 0.3, 0.95, 0.7, NAVY_DEEP, rounded=True, radius=0.12)
    txt(s, x + 0.28, iy + 0.3, 0.95, 0.7, [[(sym, {"size": 22, "bold": True, "color": ORANGE, "font": "Cambria"})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 1.4, iy + 0.34, iw - 1.6, 0.4, [[(head, {"size": 14, "bold": True, "color": NAVY})]])
    txt(s, x + 0.28, iy + 1.12, iw - 0.5, 0.5, [[(body, {"size": 12, "color": INK})]], line_spacing=1.02)
chips(s, 0.62, 6.35, ["statistical model", "matrix formulation", "interaction effects"])
footer(s, "07")

# =====================================================================
# 8 — TESTS & ASSUMPTIONS  (Carl Johan)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEP 5b · INFERENCE", "Parametric and non-parametric tests — with assumptions")
speaker_pill(s, "Carl Johan")

# two test columns
col = [("Parametric", NAVY, [
            "Normal-approximation 95% CI for the paired Δ (paired standard error).",
            "Wilson intervals on each arm's compliance rate (proportions near 0/1)."]),
       ("Exact / non-parametric", ORANGE_DK, [
            "McNemar's exact (binomial) test on the discordant pairs.",
            "Category-clustered bootstrap — 20,000 resamples (see next slides)."])]
for i, (head, c, its) in enumerate(col):
    x = 0.62 + i * (3.95 + 0.3)
    rect(s, x, 2.0, 3.95, 2.5, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0, shadow=True)
    txt(s, x + 0.28, 2.2, 3.5, 0.4, [[(head, {"size": 15.5, "bold": True, "color": c})]])
    bullets(s, x + 0.28, 2.7, 3.45, 1.7, its, size=12, gap=8, lh=1.05, marker_color=c)

# assumptions panel
ax = 9.15
rect(s, ax, 2.0, 3.57, 4.5, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, ax + 0.3, 2.2, 3.0, 0.4, [[("ASSUMPTIONS & CHECKS", {"size": 11, "bold": True, "color": ORANGE, "spacing": 1.4})]])
bullets(s, ax + 0.3, 2.66, 3.0, 3.7, [
    [("Independence", {"bold": True, "color": WHITE}), (" of prompts — only approximate, so we add a category-clustered bootstrap.", {"color": ICE})],
    [("Binary → Bernoulli", {"bold": True, "color": WHITE}), ("; deterministic decoding removes sampling noise.", {"color": ICE})],
    [("Near 0/1 rates", {"bold": True, "color": WHITE}), (" → Wilson / exact instead of a naive normal interval.", {"color": ICE})],
    [("McNemar", {"bold": True, "color": WHITE}), (" needs only the discordant pairs to be exchangeable under H₀.", {"color": ICE})],
], size=12, gap=8, lh=1.04, marker_color=ORANGE)

# bottom strip
rect(s, 0.62, 4.7, 8.2, 1.8, CARD, rounded=True, radius=0.06, line=CARDLINE, line_w=1.0)
txt(s, 0.92, 4.9, 7.6, 0.4, [[("Why both?", {"size": 14, "bold": True, "color": NAVY})]])
txt(s, 0.92, 5.3, 7.7, 1.1,
    [[("The parametric CI gives an interpretable effect ", {"size": 13, "color": INK}),
      ("size with uncertainty", {"size": 13, "bold": True, "color": INK}),
      ("; the exact test makes ", {"size": 13, "color": INK}),
      ("no large-sample / normality assumption", {"size": 13, "bold": True, "color": INK}),
      (" — important when 99% of the action is in a few hundred discordant pairs. They agree, which is the robustness story.", {"size": 13, "color": INK})]],
    line_spacing=1.08)
footer(s, "08")

# =====================================================================
# 9 — MAIN RESULT  (Carl Johan)
# =====================================================================
s = add_slide()
title_block(s, "RUBRIC STEP 5 · RESULT", "The edit removes most benchmark refusal", tw=8.5)
speaker_pill(s, "Carl Johan")

sx = 0.62
rect(s, sx, 2.05, 4.55, 1.2, WHITE, rounded=True, radius=0.07, line=LBLUE, line_w=2.0, shadow=True)
txt(s, sx + 0.3, 2.16, 4.0, 0.32, [[("A · BASELINE", {"size": 11, "bold": True, "color": NAVY, "spacing": 1.4})]])
txt(s, sx + 0.3, 2.46, 4.0, 0.7, [[("21.4%", {"size": 34, "bold": True, "color": NAVY, "font": TITLE_FONT}), ("   94/440", {"size": 13, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, sx, 3.4, 4.55, 1.2, WHITE, rounded=True, radius=0.07, line=ORANGE, line_w=2.0, shadow=True)
txt(s, sx + 0.3, 3.51, 4.0, 0.32, [[("B · EDITED", {"size": 11, "bold": True, "color": ORANGE_DK, "spacing": 1.4})]])
txt(s, sx + 0.3, 3.81, 4.0, 0.7, [[("76.6%", {"size": 34, "bold": True, "color": ORANGE_DK, "font": TITLE_FONT}), ("   337/440", {"size": 13, "color": MUTED})]], anchor=MSO_ANCHOR.MIDDLE)
rect(s, sx, 4.8, 4.55, 1.6, NAVY, rounded=True, radius=0.06, shadow=True)
txt(s, sx + 0.3, 4.94, 4.0, 0.32, [[("PAIRED DIFFERENCE", {"size": 11, "bold": True, "color": ORANGE, "spacing": 1.4})]])
txt(s, sx + 0.3, 5.22, 4.0, 0.6, [[("+55.2 pp", {"size": 31, "bold": True, "color": WHITE, "font": TITLE_FONT})]])
txt(s, sx + 0.3, 5.84, 4.1, 0.45, [[("95% CI [50.4, 60.0]   ·   McNemar p ≈ 9×10⁻⁶⁸", {"size": 11, "color": ICE})]])

add_image(s, os.path.join(FIG, "overall_compliance.png"), 5.45, 2.15, w=7.25)
txt(s, 5.45, 6.5, 7.25, 0.4,
    [[("Same model, one deterministic edit. ", {"size": 12.5, "italic": True, "bold": True, "color": INK}),
      ("Error bars: Wilson 95% intervals.", {"size": 12.5, "italic": True, "color": MUTED})]], align=PP_ALIGN.CENTER)
footer(s, "09")

# =====================================================================
# 10 — RESAMPLING & PAIRED FLIPS  (Oliver)
# =====================================================================
s = add_slide()
title_block(s, "RESAMPLING & ROBUSTNESS", "Resampling confirms a large, one-directional effect")
speaker_pill(s, "Oliver")

# 2x2 mini table
tx, ty, cell, labw = 0.62, 2.35, 1.5, 1.55
txt(s, tx + labw, ty - 0.4, cell * 2, 0.32, [[("EDITED", {"size": 10.5, "bold": True, "color": ORANGE_DK, "spacing": 1.4})]], align=PP_ALIGN.CENTER)
txt(s, tx + labw, ty, cell, 0.42, [[("Refuse", {"size": 11, "bold": True, "color": INK})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx + labw + cell, ty, cell, 0.42, [[("Comply", {"size": 11, "bold": True, "color": INK})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
rowlabs = ["Refuse", "Comply"]
data = [[("99", False), ("247", True)], [("4", True), ("90", False)]]
for r in range(2):
    ry = ty + 0.42 + r * cell
    txt(s, tx - 0.2, ry, labw, cell, [[(rowlabs[r], {"size": 11, "bold": True, "color": INK})]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    for c in range(2):
        cxp = tx + labw + c * cell
        val, key = data[r][c]
        fillc = ORANGE if (r == 0 and c == 1) else (NAVY if (r == 1 and c == 0) else WHITE)
        textc = WHITE if key else INK
        rect(s, cxp, ry, cell, cell, fillc, line=CARDLINE, line_w=1.0, shadow=key)
        txt(s, cxp, ry, cell, cell, [[(val, {"size": 30, "bold": True, "color": textc, "font": TITLE_FONT})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, tx, ty + 0.42 + 2 * cell + 0.1, labw + 2 * cell, 0.35,
    [[("Baseline (rows) → Edited (cols).  247 vs 4 discordant flips.", {"size": 10.5, "italic": True, "color": MUTED})]])

# right cards
px = 6.7
rect(s, px, 2.05, 6.02, 1.95, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, px + 0.32, 2.22, 5.4, 0.4, [[("Category-clustered bootstrap", {"size": 15, "bold": True, "color": ORANGE})]])
txt(s, px + 0.32, 2.66, 5.5, 1.2,
    [[("20,000 resamples (seed 2445), resampling whole categories. ", {"size": 13, "color": ICE}),
      ("Bootstrap 95% CI [47.3, 63.0]", {"size": 13.5, "bold": True, "color": WHITE}),
      (" — wider than the normal CI, still far from zero.", {"size": 13, "color": ICE})]],
    line_spacing=1.08)
rect(s, px, 4.2, 6.02, 2.2, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, px + 0.32, 4.4, 5.4, 0.4, [[("Robustness checks", {"size": 15, "bold": True, "color": NAVY})]])
bullets(s, px + 0.32, 4.85, 5.45, 1.5, [
    [("McNemar's exact test on the 251 discordant pairs: ", {}), ("p ≈ 9×10⁻⁶⁸", {"bold": True, "color": ORANGE_DK}), (".", {})],
    [("Sensitivity — dropping the 8 length-capped answers still gives ", {}), ("+56.0 pp", {"bold": True, "color": NAVY}), (".", {})],
    "The four opposite flips are borderline cases → error analysis, not a safety gain.",
], size=12, gap=7, lh=1.04)
chips(s, 0.62, 6.55, ["resampling / bootstrap", "uncertainty & stability"])
footer(s, "10")

# =====================================================================
# 11 — GENERALIZATION + FAIRNESS  (Oliver)
# =====================================================================
s = add_slide()
title_block(s, "GENERALIZATION & FAIRNESS LENS", "Does it generalize — and is coverage uniform?", tw=10.5)
speaker_pill(s, "Oliver")

# left: generalization with stability fig
txt(s, 0.62, 1.95, 6.1, 0.4, [[("Generalization (hold-out + repeated splits)", {"size": 14.5, "bold": True, "color": NAVY})]])
add_image(s, os.path.join(FIG, "direction_stability.png"), 0.55, 2.35, w=5.95)
txt(s, 0.62, 5.75, 6.1, 0.75,
    [[("Direction stable across held-out policy areas (median cosine 0.95–0.97; 5 random splits/budget = CV-style). SORRY-Bench is a fully external hold-out.", {"size": 11.5, "color": INK})]],
    line_spacing=1.05)

# right: fairness lens
px = 7.0
rect(s, px, 1.95, 5.72, 4.55, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, px + 0.3, 2.15, 5.1, 0.4, [[("Bias / fairness-of-coverage lens", {"size": 15, "bold": True, "color": ORANGE_DK})]])
bullets(s, px + 0.3, 2.62, 5.15, 3.7, [
    [("Treat the 44 harm categories as ", {}), ("subgroups", {"bold": True, "color": NAVY}), (" and ask whether the safety change is uniform.", {})],
    [("It is not: residual refusals ", {}), ("concentrate", {"bold": True, "color": ORANGE_DK}), (" — child-related 7/10, then PII, harassment, financial advice.", {})],
    [("So safety ", {}), ("coverage differs across harm types", {"bold": True, "color": NAVY}), (" — a fairness-style disaggregated finding (10/category → exploratory).", {})],
    [("Responsible handling: child-related content was excluded from the construction data.", {"italic": True, "color": MUTED})],
], size=12.5, gap=10, lh=1.05)
footer(s, "11")

# =====================================================================
# 12 — CORE CONCLUSIONS  (single closing slide)  (Oliver)
# =====================================================================
s = add_slide(dark=True)
ringc = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2.0), Inches(4.4), Inches(5.4), Inches(5.4))
ringc.shadow.inherit = False; ringc.fill.background(); ringc.line.color.rgb = rgb(NAVY); ringc.line.width = Pt(1.3)
txt(s, 0.9, 0.7, 9.0, 0.35, [[("CORE CONCLUSIONS", {"size": 13, "bold": True, "color": ORANGE, "spacing": 2.6})]])
txt(s, 0.9, 1.05, 11.4, 1.0, [[("One edit, measured the way an evaluation should be", {"size": 30, "bold": True, "color": WHITE, "font": TITLE_FONT})]])
speaker_pill(s, "Oliver", dark=True)

cols = [
    ("THE RESULT", ORANGE,
     "One refusal-direction edit takes Gemma from 21.4% to 76.6% compliance — Δ +55.2 pp, flips 247 vs 4, McNemar p ≈ 9×10⁻⁶⁸. Parametric, exact, and bootstrap inference agree."),
    ("THE SCOPE", LBLUE,
     "A benchmark-specific behaviour change, not a verdict on overall quality or every safety mechanism. Effect is uneven across harm categories."),
    ("THE IMPLICATION", "FFFFFF",
     "Static refusal rates don't capture robustness for open-weight models. Evaluations should add robustness-to-editing, generalization, and fresh hold-out tests."),
]
cw3, cy = 3.83, 2.3
for i, (h, c, body) in enumerate(cols):
    x = 0.9 + i * (cw3 + 0.2)
    rect(s, x, cy, cw3, 3.35, NAVY_DEEP, rounded=True, radius=0.05, shadow=True)
    txt(s, x + 0.32, cy + 0.3, cw3 - 0.6, 0.4, [[(h, {"size": 12.5, "bold": True, "color": c, "spacing": 1.6})]])
    txt(s, x + 0.32, cy + 0.82, cw3 - 0.62, 2.4, [[(body, {"size": 13, "color": ICE})]], line_spacing=1.12)
rect(s, 0.9, 6.05, 11.53, 0.78, NAVY, rounded=True, radius=0.1)
txt(s, 0.9, 6.05, 11.53, 0.78,
    [[("Code & data:  ", {"size": 13, "color": ICE}),
      ("github.com/MatthesF/refusal-ablation", {"size": 13.5, "bold": True, "color": WHITE}),
      ("        Thank you!", {"size": 13.5, "italic": True, "color": ORANGE})]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# APPENDIX (backup for individual questioning)
# =====================================================================
def appendix_header(s, kicker, title, owner):
    txt(s, 0.62, 0.5, 9.6, 0.3, [[("APPENDIX · " + kicker, {"size": 11.5, "bold": True, "color": MUTED, "spacing": 2.2})]])
    txt(s, 0.62, 0.82, 9.6, 1.0, [[(title, {"size": 27, "bold": True, "color": INK, "font": TITLE_FONT})]], line_spacing=0.98)
    speaker_pill(s, owner, sublabel="Q&A LEAD")

# A1 — GLM design matrix detail
s = add_slide()
appendix_header(s, "A1 · STATISTICAL MODEL", "Design matrix & interaction in full", "Carl Johan")
txt(s, 0.62, 2.0, 7.2, 0.5,
    [[("logit P(Yᵢ = 1) = β₀ + β₁·edited + Σ_c γ_c·catᵢ,c + Σ_c δ_c·(edited × catᵢ,c)", {"size": 16, "bold": True, "color": NAVY, "font": "Cambria"})]])
bullets(s, 0.62, 2.7, 7.2, 3.5, [
    [("Rows of X", {"bold": True, "color": NAVY}), (" = the 880 (prompt × condition) observations; columns = intercept, condition, 43 category dummies, and condition×category.", {})],
    [("β₁", {"bold": True, "color": ORANGE_DK}), (" is the main A/B treatment effect; the McNemar test and Δ are the focused, assumption-light version of this term.", {})],
    [("δ_c", {"bold": True, "color": ORANGE_DK}), (" are interaction effects: they let the edit's effect vary by harm category (exploratory — 10 prompts each).", {})],
    [("Link choice", {"bold": True, "color": NAVY}), (": logistic for a principled binary model; a linear-probability version reproduces Δ directly as β₁.", {})],
], size=13, gap=11, lh=1.05)
# mini design-matrix sketch
mx = 8.3
rect(s, mx, 2.0, 4.4, 4.2, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, mx + 0.3, 2.18, 3.8, 0.3, [[("X  (schematic)", {"size": 12.5, "bold": True, "color": NAVY})]])
hdrs = ["1", "edit", "cat₂…", "edit×cat"]
cw = 0.92
for j, hd in enumerate(hdrs):
    txt(s, mx + 0.35 + j * cw, 2.55, cw, 0.3, [[(hd, {"size": 10, "bold": True, "color": MUTED})]], align=PP_ALIGN.CENTER)
mat = [["1", "0", "…", "0"], ["1", "1", "…", "1"], ["1", "0", "…", "0"], ["1", "1", "…", "1"]]
for r, row in enumerate(mat):
    ry = 2.9 + r * 0.62
    for j, v in enumerate(row):
        cell_c = ORANGE_DK if (j in (1, 3) and v == "1") else INK
        txt(s, mx + 0.35 + j * cw, ry, cw, 0.4, [[(v, {"size": 13, "bold": (cell_c == ORANGE_DK), "color": cell_c, "font": "Cambria"})]], align=PP_ALIGN.CENTER)
txt(s, mx + 0.3, 5.55, 3.9, 0.5, [[("Paired rows share their category dummies; only the condition columns flip.", {"size": 10.5, "italic": True, "color": MUTED})]], line_spacing=1.03)
footer(s, "A1", tag="Appendix")

# A2 — test decision table
s = add_slide()
appendix_header(s, "A2 · INFERENCE", "Which test, and why", "Carl Johan")
tx, ty = 0.62, 2.0
tw_tbl = 12.1
colw = [4.0, 3.0, 2.4, 2.7]
heads = ["Question", "Test", "Type", "Key assumption"]
rect(s, tx, ty, tw_tbl, 0.55, NAVY)
cxp = tx + 0.18
for hd, w in zip(heads, colw):
    txt(s, cxp, ty, w - 0.2, 0.55, [[(hd, {"size": 12, "bold": True, "color": WHITE})]], anchor=MSO_ANCHOR.MIDDLE)
    cxp += w
rowsT = [
    ("Did the edit change compliance?", "McNemar's exact (binomial)", "Non-param / exact", "discordant pairs exchangeable"),
    ("How large, with uncertainty?", "Normal-approx CI on paired Δ", "Parametric", "large n; paired SE valid"),
    ("Robust to category clustering?", "Clustered bootstrap (20k)", "Resampling", "categories exchangeable"),
    ("Per-arm rate uncertainty?", "Wilson interval", "Parametric", "binomial proportion"),
    ("Does effect vary by category?", "Logistic GLM w/ interaction", "Parametric (exploratory)", "model & link correct"),
]
rh = 0.78
for i, row in enumerate(rowsT):
    ry = ty + 0.55 + i * rh
    rect(s, tx, ry, tw_tbl, rh, CARD if i % 2 == 0 else WHITE, line=CARDLINE, line_w=0.75)
    cxp = tx + 0.18
    styles = [(12.5, True, INK), (12, False, NAVY), (11.5, True, ORANGE_DK), (11.5, False, MUTED)]
    for (val, w, (sz, b, cc)) in zip(row, colw, styles):
        txt(s, cxp, ry, w - 0.25, rh, [[(val, {"size": sz, "bold": b, "color": cc})]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
        cxp += w
footer(s, "A2", tag="Appendix")

# A3 — power derivation
s = add_slide()
appendix_header(s, "A3 · SAMPLE SIZE", "Power & sample-size derivation", "Matthes")
rect(s, 0.62, 2.05, 6.0, 3.4, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, 0.92, 2.3, 5.4, 0.6,
    [[("n = ( σ · (z₁₋β + z₁₋α/2) / Δ )²", {"size": 18, "bold": True, "color": NAVY, "font": "Cambria"})]])
bullets(s, 0.92, 3.05, 5.45, 2.3, [
    [("σ = 1", {"bold": True, "color": INK}), (" — conservative, since Dᵢ ∈ [−1, 1].", {})],
    [("Δ = 0.15", {"bold": True, "color": INK}), (" target effect (15 pp).", {})],
    [("z₁₋β = 0.84", {"bold": True, "color": INK}), (" (80% power),  ", {}), ("z₁₋α/2 = 1.96", {"bold": True, "color": INK}), (" (α = 0.05).", {})],
    [("→ n ≈ 349", {"bold": True, "color": ORANGE_DK}), ("  needed.", {})],
], size=13, gap=10, lh=1.05)
rect(s, 6.9, 2.05, 5.82, 3.4, NAVY, rounded=True, radius=0.05, shadow=True)
txt(s, 7.2, 2.35, 5.2, 0.4, [[("WHAT WE HAD", {"size": 11, "bold": True, "color": ORANGE, "spacing": 1.8})]])
txt(s, 7.2, 2.8, 5.3, 2.4,
    [[("440 paired prompts", {"size": 20, "bold": True, "color": WHITE, "font": TITLE_FONT})],
     [("→ 88.2% power for a 15 pp change.", {"size": 14, "color": ICE})],
     [("", {"size": 8})],
     [("Observed effect = 55.2 pp", {"size": 16, "bold": True, "color": WHITE})],
     [("— far above the planning threshold, so detection was never in doubt.", {"size": 13, "color": ICE})]],
    line_spacing=1.12, space_after=3)
footer(s, "A3", tag="Appendix")

# A4 — category detail / fairness
s = add_slide()
appendix_header(s, "A4 · DISAGGREGATION", "Category-level change (subgroup view)", "Oliver")
add_image(s, os.path.join(FIG, "category_compliance_change.png"), 0.55, 2.05, w=8.0)
txt(s, 0.62, 5.62, 8.0, 0.4, [[("Edited − baseline compliance by category. Many move 0→100%; others barely change.", {"size": 11, "italic": True, "color": MUTED})]])
rect(s, 9.2, 2.05, 3.52, 4.2, CARD, rounded=True, radius=0.05, line=CARDLINE, line_w=1.0)
txt(s, 9.45, 2.25, 3.05, 0.4, [[("Reading it as fairness", {"size": 13.5, "bold": True, "color": NAVY})]])
bullets(s, 9.45, 2.7, 3.05, 3.4, [
    "103/440 prompts still refused after the edit.",
    [("Residuals concentrate: ", {}), ("child-related 7/10", {"bold": True, "color": ORANGE_DK}), (", PII, harassment, financial advice.", {})],
    "Uneven coverage across harm types = a fairness-of-coverage signal.",
    "Small per-category n (10) → exploratory, not a ranking.",
], size=11.5, gap=9, lh=1.04)
footer(s, "A4", tag="Appendix")

# ---- presenter notes (brief) ----
NOTES = {
    0: "Oskar opens. 02445 video pitch. Headline 21.4->76.6%. Frame as an A/B evaluation of an open-weight model.",
    1: "Oskar. Step 1 aspect: robustness/red-teaming of open-weight refusal. Why it matters: benchmark validity.",
    2: "Oskar. Steps 2-3: paired A/B test (A baseline vs B edited), metric = official Mistral LLM-as-a-Judge, binary 0/1.",
    3: "Matthes. Step 3 model+runtime: gemma-4-E4B-it, open-weight, RunPod RTX 5090, deterministic, independent samples.",
    4: "Matthes. Step 4 data: SORRY-Bench benchmark + 480 self-curated (12 areas x safe/unsafe), disjoint splits, power n=349 -> 88.2%.",
    5: "Matthes. Step 4c: the edit in one line; answers -> Bernoulli Y in {0,1}; paired (Y_base, Y_edit).",
    6: "Carl Johan. Model: paired Delta; GLM/matrix lens with main effect beta1 and category interaction delta_c. Hits objectives 1-3.",
    7: "Carl Johan. Parametric (normal CI, Wilson) + exact/non-parametric (McNemar) + assumptions and how checked.",
    8: "Carl Johan. Result: 21.4->76.6%, +55.2pp, CI [50.4,60.0], McNemar p~9e-68.",
    9: "Oliver. Resampling: clustered bootstrap CI [47.3,63.0]; flips 247 vs 4; sensitivity 56.0pp.",
    10: "Oliver. Generalization (held-out areas, CV-style splits, external benchmark) + fairness-of-coverage across categories.",
    11: "Oliver. Single core-conclusions slide: result / scope / implication + repo. Thank you.",
}
for idx, slide in enumerate(prs.slides):
    if idx in NOTES:
        slide.notes_slide.notes_text_frame.text = NOTES[idx]

out = os.path.join(HERE, "Refusal_Ablation_Defense_v2.pptx")
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
